#!/usr/bin/env python3
"""生成参数化意图大纲（Blueprint Compiler）原理 PPT"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 色彩定义
HW_RED    = RGBColor(0xC7, 0x00, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG   = RGBColor(0x2E, 0x2E, 0x2E)
ALT_BG    = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_RED = RGBColor(0xFF, 0xE5, 0xE5)
LIGHT_BLU = RGBColor(0xE3, 0xF2, 0xFD)
FONT_CN   = "微软雅黑"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_text(l, t, w, h, text, size, color=DARK_BG, bold=False, align=PP_ALIGN.LEFT, font=FONT_CN):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return box

def add_arrow(l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = HW_RED
    s.line.fill.background()

# ── 标题与修饰 ──
add_rect(0,       0, 33.87, 2.0, HW_RED)
add_rect(0,   18.65, 33.87, 0.4, HW_RED)
add_text(1.0, 0.2, 30, 1.0, "报告大纲架构：参数化意图 → Agent 自主编译 → 底层执行骨架", 20, color=WHITE, bold=True)
add_text(1.0, 1.2, 30, 0.8, "大纲只描述「意图的可变参数」(分析哪个指标/看什么时段/补充什么观察)，Agent 自主决定「如何展现」(图表/表格/文本)", 10, color=WHITE)

# ==========================================
# 左侧：大纲意图层
# ==========================================
add_text(1.0,   2.5, 12.0, 0.5, "第一层：参数化意图大纲 (outline)", 14, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(1.0,   3.2, 12.0, 0.8, "用户看到的是一句连贯的话，带可调参数插槽。完全不涉及任何展现形式（图表/表格）的决定。", 8.5, color=DARK_BG)
add_rect(1.0,   4.3, 12.0, 14.0, ALT_BG)

tree1 = """outline:
  # [连贯的意图描述]
  document: |
    对设备 {@target_device} 的
    {@focus_metric} 进行深度分析。
    分析范围为 {@analysis_period}。
    如有异常,请结合
    {@supplementary_context}
    给出专业诊断建议。
    
  # [意图的可变参数]
  blocks:
    - id: "target_device"
      type: "param_ref"
      default: "{$device}"
      
    - id: "focus_metric"
      type: "indicator"
      default: "振动幅值"
      
    - id: "analysis_period"
      type: "time_range"
      default: "近3天"
      
    - id: "supplementary_context"
      type: "free_text"
      default: "端盖有轻微漏油" """

# 先画参数高亮背景
add_rect(1.3,  8.7, 11.4, 1.1, LIGHT_BLU)
add_rect(1.3, 10.5, 11.4, 1.1, LIGHT_BLU)
add_rect(1.3, 12.3, 11.4, 1.1, LIGHT_BLU)
add_rect(1.3, 14.1, 11.4, 1.1, LIGHT_BLU)
# 再画代码文本
add_text(1.2, 4.4, 11.6, 13.5, tree1, 8.5, font=FONT_MONO)

# ==========================================
# 中间：编译层
# ==========================================
add_text(14.0, 7.0, 4.0, 0.6, "Agent (Template Copilot)", 9, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(14.0, 7.8, 4.0, 1.6, "通读整段意图\n自主决策：\n用什么数据源\n用什么展现形式", 8, color=DARK_BG, align=PP_ALIGN.CENTER)
add_arrow(14.5, 9.5, 3.0, 0.6)

# ==========================================
# 右侧：运行层
# ==========================================
add_text(19.0,   2.5, 13.0, 0.5, "第二层：编译产物 (content)", 14, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(19.0,   3.2, 13.0, 0.8, "Agent 自主选择的数据源和展现方式。用户没有指定任何图表类型——全是 Agent 编译产物。", 8.5, color=DARK_BG)
add_rect(19.0,   4.3, 13.0, 14.0, ALT_BG)

tree2 = """content:
  datasets:
    # Agent 自主决定需要 3 个数据源
    - id: "ds_vib_current"
      source: { kind: nl2sql }
      desc: "查设备最近一次振幅"
      
    - id: "ds_vib_trend"
      source: { kind: nl2sql }
      desc: "查近3天振幅历史"
      
    - id: "ds_diagnosis"
      source: { kind: ai_synthesis }
      prompt: "分析+漏油事实"
      
  presentation:
    # Agent 自主决定排版方案
    type: composite_table
    sections:
      - band: "诊断结论"
        layout: kv_grid
        
      - band: "趋势图"      
        layout: chart(line) """

# 高亮 Agent 自主决策部分
add_rect(19.3,   5.5, 12.4, 2.5, LIGHT_RED)
add_rect(19.3,   8.6, 12.4, 2.0, LIGHT_RED)
add_rect(19.3,  11.2, 12.4, 2.5, LIGHT_RED)
add_rect(19.3,  15.3, 12.4, 2.5, LIGHT_RED)
# 再画代码
add_text(19.2, 4.4, 12.6, 13.5, tree2, 8.5, font=FONT_MONO)

# 保存
out = r"e:\code\antigravity_projects\ReportTemplate\output\outline_compiler.pptx"
prs.save(out)
print(f"Saved: {out}")
