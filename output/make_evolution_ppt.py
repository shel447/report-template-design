#!/usr/bin/env python3
"""生成大纲交互三态演化关系对比 PPT"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 色彩定义
HW_RED    = RGBColor(0xC7, 0x00, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG   = RGBColor(0x2E, 0x2E, 0x2E)
ALT_BG    = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_RED = RGBColor(0xFF, 0xE5, 0xE5)
FONT_CN   = "微软雅黑"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_text(l, t, w, h, text, size, color=DARK_BG, bold=False, align=PP_ALIGN.LEFT, font=FONT_CN):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return box

def add_arrow(l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0xA0, 0xA0, 0xA0)
    s.line.fill.background()

# ── 标题与修饰 ──
add_rect(0,       0, 33.87, 2.0, HW_RED)
add_rect(0,   18.65, 33.87, 0.4, HW_RED)
add_text(1.0, 0.2, 30, 1.0, "大纲意图交互 (Outline WYSIWYG)：模板实例结构的三态演化", 20, color=WHITE, bold=True)
add_text(1.0, 1.2, 30, 0.8, "追踪同一个 Schema 实例在生命周期内从「纯模板」到「初稿态」再到被 Agent 动态重构的「定稿态」", 10, color=WHITE)

COL_W = 9.8
COL1_L = 1.0
COL2_L = 12.0
COL3_L = 23.0

add_arrow(11.0, 8.0, 0.8, 0.4)
add_arrow(22.0, 8.0, 0.8, 0.4)

# ==========================================
# 左塔：纯模板态
# ==========================================
add_text(COL1_L,   2.5, COL_W, 0.5, "阶段一：纯模板态", 14, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(COL1_L,   3.2, COL_W, 0.8, "技术人员刚配置好，未接入任何真实机器批次数据时的静态定义集。", 8, color=DARK_BG)
add_rect(COL1_L,   4.3, COL_W, 14.0, ALT_BG)

tree1 = """outline:
  draft_prompt: "..."


content:
  datasets:
    # [1] 原生 sql
    - id: "ds_vib"  
      source: { kind: sql }
      
    # [2] 原生 ai 分析
    - id: "ds_ai"   
      depends_on: ["ds_vib"]
      source:
        kind: ai_synthesis
        prompt: "分析数据"

  presentation:
    sections:
      # [a] 原生布局块
      - band: "诊断"
        layout: kv_grid
        dataset_id: "ds_ai" """
add_text(COL1_L+0.2, 4.4, COL_W-0.4, 13.5, tree1, 8.5, font=FONT_MONO)

# ==========================================
# 中塔：草稿态
# ==========================================
add_text(COL2_L,   2.5, COL_W, 0.5, "阶段二：草稿生成态", 14, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(COL2_L,   3.2, COL_W, 0.8, "引擎刚拉取完基础事实，拿着 draft_prompt 调用一次 LLM 吐出预览文字，呈现给用户的底本。", 8, color=DARK_BG)
add_rect(COL2_L,   4.3, COL_W, 14.0, ALT_BG)

tree2 = """outline:
  draft_prompt: "..."
  # [系统初版块列阵留档] 
  ★ original_blocks:
     - type: text
       content: "当前振
       幅8.5,存在风险。"



# 👇 以下与第一阶段完全一致
content:
  datasets:
    # [1] 原生 sql
    - id: "ds_vib"  
      ...
    # [2] 原生 ai 分析
    - id: "ds_ai"   
      ...
      
      
  presentation:
    sections:
      # [a] 原生布局块
      - band: "诊断" ..."""

# 先画背景高亮
add_rect(COL2_L+0.3, 5.1, COL_W-0.6, 2.0, LIGHT_RED)
# 再画代码
add_text(COL2_L+0.2, 4.4, COL_W-0.4, 13.5, tree2, 8.5, font=FONT_MONO)


# ==========================================
# 右塔：定稿态
# ==========================================
add_text(COL3_L,   2.5, COL_W, 0.5, "阶段三：用户定稿态 (Agent 编译结果)", 14, color=HW_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(COL3_L,   3.2, COL_W, 0.8, "用户在界面修改后，Template Copilot 据此重构出来的被篡改且扩容终极 AST。", 8, color=DARK_BG)
add_rect(COL3_L,   4.3, COL_W, 14.0, ALT_BG)

tree3 = """outline:
  original_blocks: ...
  # [Notion式混合区块]
  ★ user_blocks:
     - type: text
       content: "看重度漏油"
     ★- type: intent
       command: chart
       content: "加三日折线"

content:
  datasets:
    # [1] 原生 sql (客观基石，不动)
    - id: "ds_vib" ...
    # [2] 原生 ai 分析
    - id: "ds_ai"
      prompt: "【★外挂现场漏油事实】"
    # [3] Agent 动态响应新诉求扩种
    ★- id: "ds_trend(Agent生)"
      source: { kind: nl2sql }
      
  presentation:
    sections:
      # [a] 原生布局块的壳子保留
      - band: "诊断"
      # [b] 动态给用户搭的新台子
      ★- band: "趋势图(Agent加)"
        layout: chart
        dataset_id: "ds_trend" """

# 绘制背景高亮块
add_rect(COL3_L+0.3,  5.1, COL_W-0.6, 2.8, LIGHT_RED) # user_blocks 区域
add_rect(COL3_L+0.3, 10.6, COL_W-0.6, 0.6, LIGHT_RED) # prompt 篡改区域
add_rect(COL3_L+0.3, 11.2, COL_W-0.6, 1.2, LIGHT_RED) # nl2sql 添加区域
add_rect(COL3_L+0.3, 14.5, COL_W-0.6, 1.4, LIGHT_RED) # chart 添加区域

# 再绘制树形代码文本
add_text(COL3_L+0.2, 4.4, COL_W-0.4, 13.5, tree3, 8.5, font=FONT_MONO)

# 保存
out = r"e:\code\antigravity_projects\ReportTemplate\output\outline_evolution.pptx"
prs.save(out)
print(f"Saved: {out}")
