#!/usr/bin/env python3
"""生成报告模板结构 PPT（华为红色风格，两页，带树状结构图示）"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 色彩定义
HW_RED    = RGBColor(0xC7, 0x00, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1F, 0x1F, 0x1F)
ALT       = RGBColor(0xF5, 0xF5, 0xF5)
LT_RED    = RGBColor(0xFF, 0xCC, 0xCC)
HW_GREY   = RGBColor(0x55, 0x55, 0x55)
FONT_CN   = "微软雅黑"
FONT_MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)

# ==========================================
# 辅助绘图类
# ==========================================
class PPTBuilder:
    def __init__(self, title_text, subtitle_text):
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.rect(0, 0, 33.87, 19.05, WHITE)
        self.rect(0, 0, 33.87, 2.4, HW_RED)
        self.rect(0, 18.55, 33.87, 0.5, HW_RED)
        self.rect(14.0, 2.5, 0.06, 15.9, RGBColor(0xDD, 0xDD, 0xDD))  # 垂直分隔线（现在左侧树状结构较窄）
        
        self.tb(title_text, 1.0, 0.2, 30, 1.3, 20, bold=True, color=WHITE)
        self.tb(subtitle_text, 1.0, 1.45, 32, 0.8, 8, color=LT_RED)
        self.tb("© Report Template System v2.0  |  2026-03", 1.0, 18.55, 20, 0.48, 7, color=WHITE)

    def rect(self, l, t, w, h, fill, line=False):
        s = self.slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if not line: s.line.fill.background()
        return s

    def tb(self, text, l, t, w, h, size, bold=False, color=DARK, align=PP_ALIGN.LEFT, font=FONT_CN):
        box = self.slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        p.space_after = Pt(2)  # 轻微行距
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
        return box

    def label(self, text, l, t, w=15.0):
        self.rect(l, t+0.09, 0.18, 0.44, HW_RED)
        self.tb(text, l+0.28, t-0.03, w, 0.58, 9, bold=True, color=HW_RED)
        
    def add_tree(self, text, l, t, w=13.0, h=14.0):
        box = self.slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
        tf = box.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run(); r.text = text
        r.font.size = Pt(8.5)
        r.font.name = FONT_MONO
        r.font.color.rgb = DARK
        return box

    def table(self, data, l, t, w, h, col_fracs, fs=7.5, hdr=1):
        rows, cols = len(data), len(data[0])
        shp = self.slide.shapes.add_table(rows, cols, Cm(l), Cm(t), Cm(w), Cm(h))
        tbl = shp.table
        for i, f in enumerate(col_fracs):
            tbl.columns[i].width = Cm(w * f)
        for ri, row in enumerate(data):
            is_h = ri < hdr
            for ci, val in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.margin_left = Cm(0.12); cell.margin_right = Cm(0.06)
                cell.margin_top = Cm(0.1); cell.margin_bottom = Cm(0.1)
                tf = cell.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]
                r = p.add_run(); r.text = str(val)
                r.font.size = Pt(fs); r.font.name = FONT_CN
                r.font.bold = is_h
                if is_h:
                    r.font.color.rgb = WHITE
                    cell.fill.solid(); cell.fill.fore_color.rgb = HW_RED
                else:
                    r.font.color.rgb = DARK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ALT if ri % 2 == 0 else WHITE
        return shp

# ==========================================
# PAGE 1: 报告模板整体结构
# ==========================================
p1 = PPTBuilder("智能报告模板系统  ·  整体目录与章节结构 (v2.0)", 
               "Template Root  |  Parameters  |  Sections Tree  |  Content Model  |  v2.0 数据-展示解耦呈现")

# ── 左侧：整体树状结构 ──
p1.label("模板整体树状结构 (YAML)", 0.5, 2.7)
tree1 = """[Template Root]
 ├── id: string                 (唯一标识)
 ├── type: string               (报告类型)
 ├── scene: string              (适用场景)
 ├── name: string               (展示名称)
 ├── parameters[]: array
 │    ├── id: string            (参数占位符 {param_id})
 │    ├── label: string
 │    ├── required: boolean
 │    ├── input_type: enum      (free_text | enum | dynamic)
 │    ├── multi: boolean        (允许多值则可做foreach)
 │    ├── options[]: array
 │    └── source: string
 └── sections[]: array          (递归章节树)
      ├── title: string         (支持 {param})
      ├── foreach: object
      │    ├── param: string    (绑定的已收集多值参数)
      │    └── as: string       (迭代变量名 {$varname})
      ├── outline: object       (★ v2.0 参数化意图蓝图)
      │    ├── document: string (带 {@id} 插槽的连贯意图文本)
      │    └── blocks[]: array  (定义占位符为 threshold/operator 等参数)
      ├── subsections[]: array  (子章节嵌套)
      │    └── ...
      └── content: object       (章节叶子节点)
           ├── datasets[]       (★ v2.0 局部数据池)
           │    ├── id: string  (局部唯一)
           │    ├── depends_on: string[] (局部DAG依赖)
           │    └── source      (数据获取)
           └── presentation     (★ v2.0 视图呈现)
                ├── type: enum  (text | chart | simple_table | composite)
                └── ..."""
p1.add_tree(tree1, 0.4, 3.5)

# ── 右侧：相关定义表格 ──
RX = 14.5; RW = 18.8
p1.label("1. 参数体系与用户交互 (parameters)", RX, 2.7, RW)
p1.table([
    ["input_type",   "说明", "提取流程"],
    ["free_text",    "自由输入文本", "LLM 自动从会话提取"],
    ["enum",         "有限枚举选项 (声明 options)", "匹配/提示"],
    ["dynamic",      "动态外部来源 (声明 source)", "执行查询获取后提示"],
], RX, 3.3, RW, 1.8, [0.15, 0.45, 0.40], fs=8)

p1.label("2. 内容叶节点结构 (content v2.0解耦)", RX, 5.5, RW)
p1.tb("v2.0 将 content 节点拆分为分离且平级的两个层级，防止耦合：", RX, 6.15, RW, 0.5, 8.5)
p1.table([
    ["分离层级", "对应字段", "说明"],
    ["Model (数据获取)", "datasets[]", "基于 id + depends_on 构建局部 DAG，控制并发与依赖"],
    ["View (视图布局)", "presentation", "仅描述 UI，通过 dataset_id 绑定数据 (单图单段或表格)"],
], RX, 6.7, RW, 1.5, [0.16, 0.16, 0.68], fs=8)

p1.label("3. 基础呈现形式 (presentation.type)", RX, 8.6, RW)
p1.table([
    ["type",           "可用 source.kind", "说明"],
    ["text",           "（无）/ ai_syn", "静态段落 ({param}直填)，或挂载 AI 报告结论摘要"],
    ["value",          "sql",            "执行单值 SQL，嵌于 anchor 单值锚点语句中"],
    ["chart",          "sql / nl2sql",   "渲染折线、柱状、饼图等，自然语言可由LLM推断 chart_type"],
    ["simple_table",   "sql / nl2sql",   "简单明细列表，自动利用查询返回结果推断列名"],
    ["composite_table","(任何组合)",      "混合复杂表格！在下页详细解析（局部的独立小世界）"],
], RX, 9.2, RW, 2.6, [0.19, 0.18, 0.63], fs=8)

p1.label("4. 章节展开机制 (foreach)", RX, 12.2, RW)
p1.tb("当一个参数被声明为 `multi: true`，系统的章节可通过 `foreach` 高效按值复制：\n例如 `device_ids = [A001, A002]`，引擎会将该 chapter 克隆为两大份，内部 {$device} 变量独立。", RX, 12.8, RW, 1.5, 8.5)

# ==========================================
# PAGE 2: 复合表格数据解耦
# ==========================================
p2 = PPTBuilder("智能报告模板系统  ·  复合表格独立定义 (composite_table)", 
               "Local DAG  |  Dataset Binding  |  ai_synthesis  |  kv_grid  |  tabular  |  Cross-referencing")

# ── 左侧：复合表格树状结构 ──
p2.label("复合表格内部树状结构 (content 节点)", 0.5, 2.7)
tree2 = """[content]
 ├── datasets[]                 (局部数据池/DAG图)
 │    ├── id: string            (必须全局唯一, eg: "ds_a")
 │    ├── depends_on[]: string  (等待另一数据源完成)
 │    └── source: object
 │         ├── kind: enum       (sql | nl2sql | ai_synthesis)
 │         ├── query / description 
 │         └── context          (★ ai_synthesis 专用)
 │              ├── refs[]      (引用同表格内已就绪的 dataset)
 │              └── queries[]   (额外补充查询)
 │         └── knowledge        (★ ai_synthesis 专用 RAG)
 │              └── params      (subject, symptoms, obj)
 │         └── prompt           (LLM 指令和框架)
 │
 └── presentation: object       (视图呈现布局)
      ├── type: "composite_table"
      ├── columns: integer      (网格总列数, eg: 8)
      └── sections[]: array     (从上至下的表格区段)
           ├── band: string     (可选, 合并列的首行标题)
           ├── dataset_id: str  (★ 绑定数据池的 id)
           └── layout: object   (布局渲染引擎)
                │
                ├── [type = kv_grid] (键值对区段)
                │    ├── cols_per_row, key_span, val_span
                │    └── fields[] (key, value/col)
                │
                └── [type = tabular] (行列表明细)
                     ├── headers[] (label, span, repeat)
                     └── columns[] (field, span, repeat)"""
p2.add_tree(tree2, 0.4, 3.5, h=14.5)

# ── 右侧：相关定义表格 ──
p2.label("1. 局部数据依赖执行顺序图 (DAG 解析法则)", RX, 2.7, RW)
p2.tb("引擎读取 datasets，分析 depends_on。无依赖节点（如 SQL 1, SQL 2）线程并发执行。\nai_synthesis 拥有显式 depends_on，引擎等待前置 SQL 就绪后再执行大模型合成逻辑：", RX, 3.3, RW, 1.0, 8.5)
# 简单绘个流程模拟
p2.rect(RX+0.5, 4.4, 3.0, 0.8, ALT, line=True); p2.tb("ds_base_sql\n(并发执行)", RX+0.5, 4.4, 3.0, 0.8, 8, align=PP_ALIGN.CENTER)
p2.rect(RX+0.5, 5.4, 3.0, 0.8, ALT, line=True); p2.tb("ds_fault_sql\n(并发执行)", RX+0.5, 5.4, 3.0, 0.8, 8, align=PP_ALIGN.CENTER)
p2.tb("──▶", RX+3.8, 4.5, 1.0, 0.6, 12, color=HW_RED)
p2.tb("──▶", RX+3.8, 5.5, 1.0, 0.6, 12, color=HW_RED)
p2.rect(RX+5.2, 4.8, 4.0, 1.0, HW_RED, line=False); p2.tb("ds_ai_summary\n(等待两前置完成, 提取 refs)", RX+5.2, 4.8, 4.0, 1.0, 8, color=WHITE, align=PP_ALIGN.CENTER)
p2.tb("──▶", RX+9.5, 4.9, 1.0, 0.6, 12, color=HW_GREY)
p2.rect(RX+10.8, 4.8, 2.8, 1.0, ALT, line=True); p2.tb("presentation\n各区段按序渲染", RX+10.8, 4.8, 2.8, 1.0, 8, align=PP_ALIGN.CENTER)

p2.label("2. 知识检索三要素 (ai_synthesis knowledge)", RX, 6.7, RW)
p2.table([
    ["要素", "配置字段", "支持提取源", "说明 / 例子"],
    ["主体", "subject", "全局参数 {param}", "是谁出问题？「西门子 S7-1200」"],
    ["征兆", "symptoms", "refs 引用 (如 ds_a.name)", "表现：动态筛选出 status='异常' 的测点清单汇聚"],
    ["目标", "objective", "静态设定", "「对比行业维护规范给出排查建议指导」"],
], RX, 7.3, RW, 2.0, [0.08, 0.16, 0.28, 0.48], fs=7.5)

p2.label("3. 布局模式：键值对栅格 (kv_grid 数据填充法则)", RX, 9.8, RW)
p2.table([
    ["绑定配置", "渲染效果与用法"],
    ["无 dataset_id", "纯静态信息展示区，使用 fields[].value 直接渲染全局变量或输入文本"],
    ["dataset_id", "单行展示模式：使用 fields[].col，按列名从查询结果首行映射填入 KV 槽"],
    ["dataset_id + col未声明", "动态多行枚举模式：引擎按查询返回的 [key_col, value_col] 动态排版无穷长行"],
    ["空余槽位", "公式：(key_span + val_span) * cols_per_row = cols。不足时自动用空白补齐"],
], RX, 10.4, RW, 2.3, [0.24, 0.76], fs=8)

p2.label("4. 布局模式：表头与明细列表 (tabular 动态能力)", RX, 13.2, RW)
p2.table([
    ["绑定配置", "渲染效果与用法"],
    ["常规 headers / cols", "静态表头定义，直接提取 query 每行的值对应填列表格的 N 列空间"],
    ["~dynamic~ 映射符", "当 query 为 PIVOT 透视动态多设备结果时引入。系统自动捕捉未声明列并克隆映射"],
    ["重复跨度 rule", "headers_span 之和必须严格等于表格声明的 Columns 总宽度配置，否则报错"],
], RX, 13.8, RW, 1.8, [0.24, 0.76], fs=8)

# 保存
out = r"e:\code\antigravity_projects\ReportTemplate\output\report_template_structure.pptx"
prs.save(out)
print(f"Saved: {out}")
