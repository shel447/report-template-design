#!/usr/bin/env python3
"""生成大纲交互运行时流程图 PPT（华为红色风格）"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── 色彩定义
HW_RED    = RGBColor(0xC7, 0x00, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG   = RGBColor(0x33, 0x33, 0x33)
ALT_BG    = RGBColor(0xF9, 0xF9, 0xF9)
LT_RED    = RGBColor(0xFF, 0xEB, 0xEB)
HW_GREY   = RGBColor(0x77, 0x77, 0x77)
FONT_CN   = "微软雅黑"

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
slide = prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid()
    if fill_color:
        s.fill.fore_color.rgb = fill_color
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def add_arrow(l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Cm(l), Cm(t), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = HW_GREY
    s.line.fill.background()

def add_text(l, t, w, h, text, size, color=DARK_BG, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = FONT_CN
    return box

# ── 标题与修饰 ──
add_rect(0, 0, 33.87, 2.4, HW_RED)
add_rect(0, 18.55, 33.87, 0.5, HW_RED)
add_text(1.0, 0.3, 30, 1.3, "报告生成运行时：大纲交与意图融合(WYSIWYG)流程", 22, color=WHITE, bold=True)
add_text(1.0, 1.45, 30, 0.8, "展示用户意图修改如何影响最终生成，以及对“越界意图”的柔性处理机制", 10, color=WHITE)

# ==========================================
# 绘制四大核心阶段
# ==========================================

configs = [
    {
        "y": 3.5,
        "step": "Phase 1: 用户编写意图大纲 (Outline Blueprint)",
        "action": "用户在富文本编辑器中直接书写自然流畅的报告意图描述。通过 '/' 唤出组件，插入例如 threshold（阈值）、operator（操作符）、indicator（指标）等受控的参数插槽。",
        "template_use": "将意图和插槽编织为\noutline.document 与 outline.blocks",
        "output": "【参数化意图文本】\n例如：“分析范围为近3天，关注存在大于(operator) 8.5(threshold) 的异常数据”",
        "color": WHITE
    },
    {
        "y": 7.0,
        "step": "Phase 2: Agent意图降维编译 (Template Copilot)",
        "action": "用户点击保存后，Agent (大模型) 通读大纲意图中的所有参数组合，自主进行结构决策：应该用什么数据源？应该用图表（并标注红线）还是表格展示？",
        "template_use": "读取 outline.document\n翻译出对应的 datasets 与 presentation",
        "output": "【编译后的布局与数据方案】\n例如：自动生成取近3天趋势的 nl2sql + 挂载折线图 chart 控件。",
        "color": ALT_BG
    },
    {
        "y": 10.5,
        "step": "Phase 3: 框架固化为标准模板",
        "action": "Agent 编译生成的技术结构被正式保存为底层标准 YAML 配置，真正做到“大纲归大纲参数，底层归严谨取数”，避免每次跑批都调用大模型的长文本推理成本。",
        "template_use": "固化到 content 节点的\ndatasets 与 presentation",
        "output": "【静态化的数据流管道】\n不再有“大纲意图”，全部降维成确定的数据流管道模型。",
        "color": LT_RED
    },
    {
        "y": 14.0,
        "step": "Phase 4: 运行时数据跑批 (Data Engine)",
        "action": "每月/每季度自动跑批时，数据引擎拉取上述固化好的模板配置，从数据库查询出真实的传感器客观数值，按照规定好的图表/表格要求完成自动渲染编织。",
        "template_use": "执行依赖尾部的 datasets \n结合真实数据渲染 presentation 视图",
        "output": "【最终正式报告】\n“生成一份包含近 3 天振幅的高精度折线图，并在 8.5 处标注了红线。”",
        "color": ALT_BG
    }
]

for i, cfg in enumerate(configs):
    y = cfg["y"]
    # 绘制背景条块
    add_rect(1.0, y, 31.87, 2.5, cfg["color"], line_color=RGBColor(210,210,210))
    
    # 左侧：阶段名
    add_rect(1.0, y, 6.0, 2.5, HW_RED)
    add_text(1.2, y + 0.9, 5.6, 1.0, cfg["step"], 12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    
    # 分割线 1
    add_rect(14.0, y, 0.05, 2.5, RGBColor(220,220,220))
    # 分割线 2
    add_rect(23.0, y, 0.05, 2.5, RGBColor(220,220,220))
    
    # 栏目 1：引擎动作
    add_text(7.5, y + 0.2, 5.0, 0.5, "▶ 引擎核心动作", 11, color=HW_RED, bold=True)
    add_text(7.5, y + 0.9, 6.0, 1.5, cfg["action"], 9, color=DARK_BG)
    
    # 栏目 2：模板触点
    add_text(14.5, y + 0.2, 5.0, 0.5, "▶ 模板信息运用 (Template)", 11, color=HW_RED, bold=True)
    add_text(14.5, y + 0.9, 8.0, 1.5, cfg["template_use"], 9, color=DARK_BG, bold=True)
    
    # 栏目 3：阶段产出 & 越界影响
    add_text(23.5, y + 0.2, 5.0, 0.5, "▶ 阶段产出与对终稿的影响", 11, color=HW_RED, bold=True)
    add_text(23.5, y + 0.9, 9.0, 1.5, cfg["output"], 9, color=DARK_BG)
    
    # 绘制向下的箭头 (除了最后一个)
    if i < len(configs) - 1:
        add_arrow(16.5, y + 2.6, 0.6, 0.8)

# 保存
out = r"e:\code\antigravity_projects\ReportTemplate\output\outline_workflow.pptx"
prs.save(out)
print(f"Saved: {out}")
